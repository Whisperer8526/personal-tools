from pptx import Presentation
from pptx.chart.data import ChartData
import pandas as pd


def replace_text_items(presentation, data_dict):
    '''
    Matches dictionary keys with encoded text placeholders in the presentation template and swaps 
    them with dictionary values. It keeps original format of the text
    Arguments:
        presentation : presentation instance
        data_dict : dictionary from first tab of Publisher-Sponsor-1
    '''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for key in data_dict.keys():
                        if key in run.text:
                            cur_text = run.text
                            new_text = cur_text.replace(key, str(data_dict[key]))
                            run.text = new_text
                            print(f"{key} --> ", run.text)

                            
def replace_table_items(presentation, data_dict):
    '''
    Matches dictionary keys with encoded table placeholders in the presentation template and swaps them
    with dictionary values. It doesn't keep original format of the table's values.
    Arguments:
        presentation : presentation instance
        data_dict : dictionary from first tab of Publisher-Sponsor-1
    ''' 
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for cell in shape.table.iter_cells():
                for key in data_dict.keys():
                    if key in cell.text:
                        cur_text = cell.text
                        new_text = cur_text.replace(key, str(data_dict[key]))
                        cell.text = new_text
                        print(f"{key} --> ", cell.text)
